#!/usr/bin/env python3
"""
H5P SlideDeck Generator

Generates H5P slide deck presentations from PDF slides and audio files
using the H5P.SlideDeck library format.
"""

import os
import json
import shutil
import zipfile
import click
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from tqdm import tqdm
from PIL import Image
import uuid
import re
from docx import Document
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation
import tempfile
import zipfile as zip

# Import project configuration
from config import PROJECT_NAME, OUTPUT_DIR, VO_DIR, PNG_DIR, NOTES_DOCX, PROJECT_TITLE, PDF_DOC


class H5PSlideDeckGenerator:
    """Main class for generating H5P SlideDeck presentations"""
    
    def __init__(self, project_name: str = None, verbose: bool = False):
        """Initialize the generator with project configuration"""
        self.project_name = project_name or PROJECT_NAME
        self.verbose = verbose
        
        # Get project root directory (where this script is located)
        # SD_Generator.py is in the project root, so we use its directory
        self.project_root = Path(os.path.dirname(os.path.abspath(__file__)))
        
        # Update paths based on project root
        self.output_dir = self.project_root / '00_Output' / self.project_name
        self.pdf_dir = self.output_dir / 'pdf_slides'  # New directory for individual PDF slides
        self.vo_dir = self.output_dir / 'VO'
        self.notes_docx = self.output_dir / f"{self.project_name}_NOTES.docx"
        self.template_dir = self.project_root / 'Template_SD'
        self.temp_dir = self.project_root / 'temp_h5p_sd_build'
        self.source_pdf = None  # Will be set when processing files
        self.slide_notes = {}  # Initialize slide notes dictionary
        
        if self.verbose:
            click.echo(f"Project: {self.project_name}")
            click.echo(f"Project Root: {self.project_root}")
            click.echo(f"PDF Directory: {self.pdf_dir}")
            click.echo(f"Audio Directory: {self.vo_dir}")
            click.echo(f"Template: SlideDeck")
    
    def validate_directories(self) -> bool:
        """Validate that all required directories exist"""
        if not self.source_pdf or not self.source_pdf.exists():
            click.echo(click.style(f"Error: Source PDF not found: {self.source_pdf}", fg='red'))
            return False
        
        if not self.vo_dir.exists():
            click.echo(click.style(f"Error: Audio directory not found: {self.vo_dir}", fg='red'))
            return False
        
        if not self.template_dir.exists():
            click.echo(click.style(f"Error: Template directory not found: {self.template_dir}", fg='red'))
            return False
        
        return True
    
    def parse_notes_document(self) -> Dict[int, str]:
        """Parse the Word document to extract slide notes"""
        slide_notes = {}
        
        if not self.notes_docx.exists():
            if self.verbose:
                click.echo(click.style(f"Warning: Notes document not found: {self.notes_docx}", fg='yellow'))
            return slide_notes
        
        try:
            doc = Document(self.notes_docx)
            current_slide = None
            current_text = []
            
            for paragraph in doc.paragraphs:
                # Check if this is a heading (slide number)
                if paragraph.style.name.startswith('Heading'):
                    # Save previous slide's text if we have one
                    if current_slide is not None and current_text:
                        slide_notes[current_slide] = ' '.join(current_text).strip()
                    
                    # Extract slide number from heading
                    heading_text = paragraph.text.strip()
                    slide_match = re.search(r'[Ss]lide\s*(\d+)', heading_text)
                    if slide_match:
                        current_slide = int(slide_match.group(1))
                        current_text = []
                    else:
                        current_slide = None
                        current_text = []
                
                elif current_slide is not None and paragraph.text.strip():
                    # This is content for the current slide
                    current_text.append(paragraph.text.strip())
            
            # Don't forget the last slide
            if current_slide is not None and current_text:
                slide_notes[current_slide] = ' '.join(current_text).strip()
            
            if self.verbose:
                click.echo(f"Parsed notes for {len(slide_notes)} slides from Word document")
                
        except Exception as e:
            click.echo(click.style(f"Error parsing notes document: {e}", fg='red'))
        
        return slide_notes

    def split_pdf_into_slides(self) -> bool:
        """Split the source PDF into individual slide PDFs"""
        try:
            # Create PDF slides directory
            self.pdf_dir.mkdir(parents=True, exist_ok=True)
            
            # Read the source PDF
            pdf = PdfReader(self.source_pdf)
            total_pages = len(pdf.pages)
            
            click.echo(f"Splitting PDF into {total_pages} slides...")
            
            # Split each page into a separate PDF
            for page_num in range(total_pages):
                writer = PdfWriter()
                writer.add_page(pdf.pages[page_num])
                
                # Save individual slide PDF
                output_path = self.pdf_dir / f"Slide{page_num + 1}.pdf"
                with open(output_path, 'wb') as output_file:
                    writer.write(output_file)
            
            click.echo(click.style("✓ PDF splitting completed successfully", fg='green'))
            return True
            
        except Exception as e:
            click.echo(click.style(f"Error splitting PDF: {e}", fg='red'))
            return False

    def get_slide_files(self) -> List[Tuple[int, Path, Optional[Path]]]:
        """Get all slide PDF files and their corresponding audio files"""
        slides = []
        
        # Get all PDF files
        pdf_files = sorted(self.pdf_dir.glob("Slide*.pdf"))
        
        for pdf_file in pdf_files:
            # Extract slide number
            match = re.match(r'Slide(\d+)\.pdf', pdf_file.name)
            if match:
                slide_num = int(match.group(1))
                
                # Find corresponding audio file
                audio_file = self.vo_dir / f"media_{slide_num}.mp3"
                if not audio_file.exists():
                    click.echo(click.style(f"Warning: No audio file for slide {slide_num}", fg='yellow'))
                    audio_file = None
                
                slides.append((slide_num, pdf_file, audio_file))
        
        return sorted(slides, key=lambda x: x[0])
    
    def generate_unique_filename(self, extension: str) -> str:
        """Generate a unique filename for H5P content"""
        # H5P uses a specific format for filenames
        unique_id = ''.join(str(uuid.uuid4()).split('-')[1:3]).upper()[:8]
        return f"{extension}-{unique_id}"
    
    def create_pdf_element(self, pdf_path: str) -> Dict:
        """Create a PDF element for the slide"""
        element = {
            "params": {
                "displayOptions": {
                    "width": "100%",
                    "height": "auto"
                },
                "pdfFile": {
                    "path": pdf_path,
                    "mime": "application/pdf",
                    "copyright": {"license": "U"}
                }
            },
            "library": "H5P.PDFViewer 1.0",
            "metadata": {
                "contentType": "PDF Viewer",
                "license": "U",
                "title": "Untitled PDF Viewer",
                "authors": [],
                "changes": [],
                "extraTitle": "Untitled PDF Viewer"
            },
            "subContentId": str(uuid.uuid4())
        }
        
        return element
    
    def create_audio_element(self, audio_path: str) -> Dict:
        """Create an audio element for the slide"""
        element = {
            "params": {
                "playerMode": "minimalistic",
                "fitToWrapper": False,
                "controls": True,
                "autoplay": True,
                "playAudio": "Play audio",
                "pauseAudio": "Pause audio",
                "contentName": "Audio",
                "audioNotSupported": "Your browser does not support this audio",
                "files": [{
                    "path": audio_path,
                    "mime": "audio/mpeg",
                    "copyright": {"license": "U"}
                }]
            },
            "library": "H5P.Audio 1.5",
            "metadata": {
                "contentType": "Audio",
                "license": "U",
                "title": "Untitled Audio",
                "authors": [],
                "changes": [],
                "extraTitle": "Untitled Audio"
            },
            "subContentId": str(uuid.uuid4())
        }
        
        return element
    
    def create_slide(self, slide_num: int, pdf_path: str, audio_path: Optional[str], 
                    notes_text: str = "") -> Dict:
        """Create a slide object for the SlideDeck presentation"""
        
        slide = {
            "image": self.create_pdf_element(pdf_path),
            "title": f"Slide {slide_num}",
            "notes": f"<p>{notes_text}</p>" if notes_text else ""
        }
        
        # Add audio element if available
        if audio_path:
            slide["audioOrVideo"] = self.create_audio_element(audio_path)
        
        return slide
    
    def generate_content_json(self, slides_data: List[Dict]) -> Dict:
        """Generate the complete content.json structure for SlideDeck"""
        content = {
            "slides": slides_data,
            "behaviour": {
                "autoPlay": False,
                "loopSlides": False
            }
        }
        
        return content
    
    def build_h5p_package(self, output_filename: str, slide_notes: Dict[int, str]) -> bool:
        """Build the H5P package with all slides and content"""
        try:
            # Debug: Print template directory path
            click.echo(f"Template directory: {self.template_dir}")
            click.echo(f"Template directory exists: {self.template_dir.exists()}")
            
            # Check if template directory exists
            if not self.template_dir.exists():
                click.echo(click.style(f"Error: Template directory not found: {self.template_dir}", fg='red'))
                return False
            
            # Create temporary directory for building
            self.temp_dir.mkdir(parents=True, exist_ok=True)
            
            # Copy entire template directory structure
            click.echo("Copying SlideDeck template files...")
            
            # Copy all template files and directories
            for item in self.template_dir.iterdir():
                if item.is_file():
                    # Copy files in root (like h5p.json)
                    shutil.copy2(str(item), str(self.temp_dir / item.name))
                elif item.is_dir() and item.name != 'content':
                    # Copy library directories (H5P.*, H5PEditor.*, FontAwesome-4.5)
                    dest_dir = self.temp_dir / item.name
                    if dest_dir.exists():
                        shutil.rmtree(str(dest_dir))
                    shutil.copytree(str(item), str(dest_dir))
            
            # Create content directory structure
            content_dir = self.temp_dir / 'content'
            content_dir.mkdir(exist_ok=True)
            
            # Create subdirectories for files and audios
            files_dir = content_dir / 'files'
            files_dir.mkdir(exist_ok=True)
            audios_dir = content_dir / 'audios'
            audios_dir.mkdir(exist_ok=True)
            
            # Get all slides
            slides = self.get_slide_files()
            if not slides:
                click.echo(click.style("Error: No slides found", fg='red'))
                return False
            
            click.echo(f"Found {len(slides)} slides")
            
            # Create slides data
            slides_data = []
            
            # Process each slide
            for slide_num, pdf_path, audio_path in tqdm(slides, desc="Processing slides"):
                # Generate unique filenames
                pdf_filename = f"pdf-{self.generate_unique_filename('pdf')}.pdf"
                audio_filename = f"audio-{self.generate_unique_filename('audio')}.mp3" if audio_path else None
                
                # Copy PDF file to content/files
                shutil.copy2(pdf_path, files_dir / pdf_filename)
                
                # Copy audio file if exists
                if audio_path:
                    shutil.copy2(audio_path, audios_dir / audio_filename)
                
                # Get notes for this slide
                notes_text = slide_notes.get(slide_num, "")
                
                # Create slide object with relative paths
                slide = self.create_slide(
                    slide_num=slide_num,
                    pdf_path=f"files/{pdf_filename}",
                    audio_path=f"audios/{audio_filename}" if audio_filename else None,
                    notes_text=notes_text
                )
                
                slides_data.append(slide)
            
            # Generate content.json
            content_data = self.generate_content_json(slides_data)
            
            # Save content.json
            content_path = content_dir / 'content.json'
            with open(content_path, 'w', encoding='utf-8') as f:
                json.dump(content_data, f, indent=2)
            
            # Update h5p.json with project title
            h5p_json_path = self.temp_dir / 'h5p.json'
            if h5p_json_path.exists():
                with open(h5p_json_path, 'r', encoding='utf-8') as f:
                    h5p_data = json.load(f)
                
                # Update title if project_title is set
                if hasattr(self, 'project_title') and self.project_title:
                    h5p_data['title'] = self.project_title
                    h5p_data['extraTitle'] = self.project_title
                else:
                    h5p_data['title'] = self.project_name
                    h5p_data['extraTitle'] = self.project_name
                
                with open(h5p_json_path, 'w', encoding='utf-8') as f:
                    json.dump(h5p_data, f, separators=(',', ':'))
            
            # Create H5P package
            click.echo(f"Creating H5P SlideDeck package: {output_filename}")
            self.output_path = self.project_root / output_filename
            
            with zip.ZipFile(self.output_path, 'w', zip.ZIP_DEFLATED) as zipf:
                for root, dirs, files in os.walk(self.temp_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, self.temp_dir)
                        zipf.write(file_path, arcname)
            
            # Clean up temporary directory
            shutil.rmtree(self.temp_dir)
            
            # Print success message with package details
            file_size = self.output_path.stat().st_size / (1024 * 1024)  # Convert to MB
            click.echo(click.style(f"✓ H5P SlideDeck package created successfully: {output_filename}", fg='green'))
            click.echo(f"  Total slides: {len(slides)}")
            click.echo(f"  File size: {file_size:.2f} MB")
            click.echo(f"  Format: SlideDeck")
            
            return True
            
        except Exception as e:
            click.echo(click.style(f"Error building H5P package: {e}", fg='red'))
            return False

    def extract_audio_from_pptx(self, pptx_path: Path) -> Tuple[bool, Dict[int, str]]:
        """Extract audio files and notes from PPTX presentation"""
        slide_notes = {}
        try:
            # Create VO directory if it doesn't exist
            self.vo_dir.mkdir(parents=True, exist_ok=True)
            
            # Open the PPTX file
            prs = Presentation(pptx_path)
            
            # Extract notes from each slide
            for idx, slide in enumerate(prs.slides, 1):
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = ""
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text"):
                            notes_text += shape.text + "\n"
                    if notes_text.strip():
                        slide_notes[idx] = notes_text.strip()
            
            # Open the PPTX file as a ZIP archive for audio extraction
            with zip.ZipFile(pptx_path, 'r') as pptx_zip:
                # Get all media files
                media_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/media/')]
                
                if not media_files:
                    click.echo(click.style("Warning: No media files found in PPTX", fg='yellow'))
                else:
                    click.echo(f"Found {len(media_files)} media files in PPTX")
                    
                    # Extract and process each media file
                    for media_file in media_files:
                        # Extract the file to a temporary location
                        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                            temp_file.write(pptx_zip.read(media_file))
                            temp_path = temp_file.name
                        
                        # Check if it's an audio file
                        if media_file.lower().endswith(('.mp3', '.wav', '.m4a')):
                            # Get slide number from filename (assuming format like media1.mp3)
                            match = re.search(r'media(\d+)', media_file)
                            if match:
                                slide_num = int(match.group(1))
                                # Copy to VO directory with standardized name
                                output_path = self.vo_dir / f"media_{slide_num}.mp3"
                                shutil.copy2(temp_path, output_path)
                                if self.verbose:
                                    click.echo(f"Extracted audio for slide {slide_num}")
                        
                        # Clean up temporary file
                        os.unlink(temp_path)
            
            if slide_notes:
                click.echo(f"Extracted notes from {len(slide_notes)} slides")
            
            # Store slide notes in the instance
            self.slide_notes = slide_notes
            
            click.echo(click.style("✓ Audio extraction completed successfully", fg='green'))
            
            return True, slide_notes
            
        except Exception as e:
            click.echo(click.style(f"Error extracting from PPTX: {e}", fg='red'))
            return False, {}


@click.command()
@click.option('--project', '-p', help='Project name (overrides config.py)')
@click.option('--output', '-o', default=None, help='Output filename (default: PROJECT_NAME_SlideDeck.h5p)')
@click.option('--pptx', '-x', help='Path to PPTX file containing audio and notes')
@click.option('--verbose', '-v', is_flag=True, help='Enable verbose output')
def main(project, output, pptx, verbose):
    """Generate H5P SlideDeck presentation from PDF and audio files"""
    
    # Initialize generator
    generator = H5PSlideDeckGenerator(project_name=project, verbose=verbose)
    
    # Extract audio and notes from PPTX if provided
    slide_notes = {}
    if pptx:
        pptx_path = Path(pptx)
        if not pptx_path.exists():
            click.echo(click.style(f"Error: PPTX file not found: {pptx_path}", fg='red'))
            return
        extracted, slide_notes = generator.extract_audio_from_pptx(pptx_path)
        if not extracted:
            return
    
    # Validate directories
    if not generator.validate_directories():
        return
    
    # Set output filename
    if not output:
        output = f"{generator.project_name}_SlideDeck.h5p"
    
    # Build the package
    generator.build_h5p_package(output, slide_notes)


if __name__ == '__main__':
    main() 